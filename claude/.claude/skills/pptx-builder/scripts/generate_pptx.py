#!/usr/bin/env python3
"""
pptx-builder: PowerPoint スライド生成スクリプト

python-pptx を使用して、JSON入力からPPTXファイルを生成する。
3つのデザインパターン（soft_business / clean_pro / warm_formal）に対応。

使い方:
    python generate_pptx.py --input slides.json --output output.pptx [--pattern clean_pro]
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
except ImportError:
    print("エラー: python-pptx がインストールされていません。")
    print("  pip install python-pptx")
    sys.exit(1)


# ============================================================
# カラーパレット
# ============================================================

COLORS = {
    "primary": RGBColor(0xC0, 0x7A, 0x5A),      # テラコッタブラウン
    "secondary": RGBColor(0x8F, 0xBA, 0x78),     # ナチュラルグリーン
    "accent": RGBColor(0x6B, 0x8F, 0xA4),        # スレートブルー
    "warm": RGBColor(0xF2, 0xDE, 0xD8),          # ピンクベージュ
    "text": RGBColor(0x3B, 0x35, 0x33),          # ダークブラウン
    "bg": RGBColor(0xF9, 0xF8, 0xF7),            # ウォームグレー
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "border": RGBColor(0xE4, 0xDF, 0xDC),
    "subtext": RGBColor(0x9A, 0x8E, 0x8B),
    "light_green": RGBColor(0xE8, 0xF0, 0xE2),
    "light_blue": RGBColor(0xE2, 0xEA, 0xF0),
    "grid": RGBColor(0xEA, 0xE5, 0xE2),
}


def C(name):
    """カラートークン名から RGBColor を取得する"""
    return COLORS[name]


# ============================================================
# デザインパターン定義
# ============================================================

PATTERNS = {
    "soft_business": {
        "heading_color": "primary",
        "slide_bg": "bg",
        "title_slide_bg": "bg",
        "header_style": "bottom_line",       # 白背景 + 下線
        "card_bg": "white",
        "card_border": "border",
        "card_radius": 0.05,
        "section_bg": "bg",
        "section_accent": None,
        "table_header_bg": "warm",
        "table_header_text": "text",
        "cta_bg": "primary",
        "cta_text": "white",
        "cta_radius": 0.04,
        "closing_bg": "bg",
    },
    "clean_pro": {
        "heading_color": "text",
        "slide_bg": "bg",
        "title_slide_bg": "bg",
        "header_style": "top_line",          # bg背景 + 上部primaryライン
        "card_bg": "white",
        "card_border": "border",
        "card_radius": 0.02,
        "section_bg": "bg",
        "section_accent": "primary",         # 左に縦線
        "table_header_bg": "primary",
        "table_header_text": "white",
        "cta_bg": "text",
        "cta_text": "white",
        "cta_radius": 0.02,
        "closing_bg": "bg",
    },
    "warm_formal": {
        "heading_color": "primary",
        "slide_bg": "bg",
        "title_slide_bg": "warm",
        "header_style": "filled",            # primary背景 + 白文字
        "card_bg": "warm",
        "card_border": None,
        "card_radius": 0.05,
        "section_bg": "warm",
        "section_accent": None,
        "table_header_bg": "primary",
        "table_header_text": "white",
        "cta_bg": "primary",
        "cta_text": "white",
        "cta_radius": 0.15,                  # ピル型
        "closing_bg": "warm",
    },
}

# レガシーテンプレート名の互換マッピング
TEMPLATE_ALIASES = {
    "minimal": "clean_pro",
    "internal": "clean_pro",
    "external": "soft_business",
}

DEFAULT_PATTERN = "clean_pro"


def resolve_pattern(name):
    """パターン名を解決する（レガシー名にも対応）"""
    name = TEMPLATE_ALIASES.get(name, name)
    if name in PATTERNS:
        return name, PATTERNS[name]
    return DEFAULT_PATTERN, PATTERNS[DEFAULT_PATTERN]


# ============================================================
# タイポグラフィ
# ============================================================

FONT_NAME = "Noto Sans JP"
FONT_FALLBACK = "Segoe UI"

SZ_TITLE = Pt(28)
SZ_SUBTITLE = Pt(16)
SZ_SECTION_HEADING = Pt(22)
SZ_SUBHEADING = Pt(16)
SZ_BODY = Pt(14)
SZ_CAPTION = Pt(10)
SZ_KPI = Pt(40)
SZ_CARD_LABEL = Pt(12)
SZ_CARD_DESC = Pt(10)
SZ_TABLE_CELL = Pt(12)
SZ_COLUMN_HEADER = Pt(16)

# ============================================================
# スライドサイズ・マージン
# ============================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LR = Inches(0.67)
MARGIN_TOP = Inches(0.55)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LR * 2
CONTENT_TOP = Inches(1.7)     # ヘッダー領域下
CONTENT_HEIGHT = Inches(5.3)


# ============================================================
# ユーティリティ関数
# ============================================================

def set_font(run, size=SZ_BODY, color_token="text", bold=False):
    """run のフォントスタイルを設定する"""
    rgb = C(color_token) if isinstance(color_token, str) else color_token
    run.font.name = FONT_NAME
    run.font.size = size
    run.font.color.rgb = rgb
    run.font.bold = bold
    # East Asian フォント
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT_NAME)
    # Latin フォールバック
    latin = rPr.find(qn("a:latin"))
    if latin is None:
        latin = rPr.makeelement(qn("a:latin"), {})
        rPr.append(latin)
    latin.set("typeface", FONT_FALLBACK)


def set_slide_bg(slide, color_token):
    """スライド背景色を設定する"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C(color_token)


def add_run(paragraph, text, size=SZ_BODY, color_token="text", bold=False):
    """段落に run を追加してフォントを設定する"""
    run = paragraph.add_run()
    run.text = text
    set_font(run, size=size, color_token=color_token, bold=bold)
    return run


def make_textbox(slide, left, top, width, height):
    """テキストボックスを追加し text_frame を返す"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def add_paragraph(tf, text, size=SZ_BODY, color_token="text",
                  bold=False, align=PP_ALIGN.LEFT, space_after=Pt(6)):
    """text_frame に段落を追加する"""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = space_after
    add_run(p, text, size=size, color_token=color_token, bold=bold)
    return p


def first_paragraph(tf, text, size=SZ_BODY, color_token="text",
                    bold=False, align=PP_ALIGN.LEFT, space_after=Pt(6)):
    """text_frame の最初の段落にテキストを設定する"""
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = space_after
    add_run(p, text, size=size, color_token=color_token, bold=bold)
    return p


def add_rect(slide, left, top, width, height, fill_token, border_token=None):
    """塗りつぶし矩形を追加する"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill_token)
    if border_token:
        shape.line.color.rgb = C(border_token)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_token,
                     border_token=None, radius=0.05):
    """角丸矩形を追加する"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill_token)
    if border_token:
        shape.line.color.rgb = C(border_token)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if shape.adjustments:
        shape.adjustments[0] = radius
    return shape


def add_line_h(slide, left, top, width, fill_token, thickness=Pt(2)):
    """水平ラインを追加する"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill_token)
    shape.line.fill.background()
    return shape


def add_line_v(slide, left, top, height, fill_token, thickness=Pt(3)):
    """垂直ラインを追加する"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, thickness, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C(fill_token)
    shape.line.fill.background()
    return shape


# ============================================================
# スライドヘッダー（コンテンツスライド共通）
# ============================================================

def add_slide_header(slide, title_text, pat):
    """パターンに応じたスライドヘッダーを描画する"""
    style = pat["header_style"]
    heading_color = pat["heading_color"]

    if style == "top_line":
        # Clean Pro: 上部に primary ライン
        add_line_h(slide, Inches(0), Inches(0), SLIDE_WIDTH, "primary", Pt(3))
        tf = make_textbox(slide, MARGIN_LR, Inches(0.35), CONTENT_WIDTH, Inches(1.0))
        first_paragraph(tf, title_text, size=SZ_SECTION_HEADING,
                        color_token=heading_color, bold=True)

    elif style == "filled":
        # Warm Formal: primary 帯 + 白文字
        add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.3), "primary")
        tf = make_textbox(slide, MARGIN_LR, Inches(0.3), CONTENT_WIDTH, Inches(0.8))
        first_paragraph(tf, title_text, size=SZ_SECTION_HEADING,
                        color_token="white", bold=True)

    else:
        # Soft Business (default): タイトル + 下線
        tf = make_textbox(slide, MARGIN_LR, Inches(0.35), CONTENT_WIDTH, Inches(1.0))
        first_paragraph(tf, title_text, size=SZ_SECTION_HEADING,
                        color_token=heading_color, bold=True)
        add_line_h(slide, MARGIN_LR, Inches(1.25), CONTENT_WIDTH, "border", Pt(1))


# ============================================================
# レイアウトハンドラー
# ============================================================

def create_title_slide(prs, data, pat):
    """4-1. タイトルスライド（表紙）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["title_slide_bg"])

    title = data.get("title", "")
    subtitle = data.get("subtitle", "")
    author = data.get("author", "")
    date = data.get("date", "")
    heading_color = pat["heading_color"]

    # タイトル
    tf = make_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.333), Inches(1.5))
    first_paragraph(tf, title, size=SZ_TITLE, color_token=heading_color,
                    bold=True, align=PP_ALIGN.CENTER)

    # 区切り線（primary、中央短い横線）
    line_w = Inches(0.7)
    line_left = (SLIDE_WIDTH - line_w) / 2
    add_line_h(slide, line_left, Inches(3.75), line_w, "primary", Pt(2))

    # サブタイトル
    if subtitle:
        tf_sub = make_textbox(slide, Inches(1.5), Inches(4.0),
                              Inches(10.333), Inches(0.8))
        first_paragraph(tf_sub, subtitle, size=SZ_SUBTITLE,
                        color_token="subtext", align=PP_ALIGN.CENTER)

    # 日付・作成者
    footer_parts = [p for p in [date, author] if p]
    if footer_parts:
        footer_text = "  |  ".join(footer_parts)
        tf_foot = make_textbox(slide, Inches(1.5), Inches(6.2),
                               Inches(10.333), Inches(0.5))
        first_paragraph(tf_foot, footer_text, size=SZ_CAPTION,
                        color_token="subtext", align=PP_ALIGN.CENTER)

    return slide


def create_toc_slide(prs, data, pat):
    """4-2. 目次スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    # 見出し
    tf_title = make_textbox(slide, MARGIN_LR, Inches(0.5), CONTENT_WIDTH, Inches(0.8))
    first_paragraph(tf_title, data.get("title", "目次"), size=SZ_SECTION_HEADING,
                    color_token="primary", bold=True)
    add_line_h(slide, MARGIN_LR, Inches(1.2), Inches(1.5), "primary", Pt(2))

    # 目次項目
    items = data.get("items", [])
    start_y = Inches(1.7)
    row_height = Inches(0.55)

    for i, item in enumerate(items):
        y = start_y + row_height * i
        num = item.get("number", f"{i + 1:02d}")
        title = item.get("title", "")
        page = item.get("page", "")

        # 番号
        tf_num = make_textbox(slide, MARGIN_LR + Inches(0.3), y,
                              Inches(0.6), Inches(0.4))
        first_paragraph(tf_num, str(num), size=SZ_SUBHEADING,
                        color_token="primary", bold=True)

        # タイトル
        tf_item = make_textbox(slide, MARGIN_LR + Inches(1.0), y,
                               Inches(8.0), Inches(0.4))
        first_paragraph(tf_item, title, size=SZ_BODY, color_token="text")

        # リーダー線
        leader_left = MARGIN_LR + Inches(9.2)
        add_line_h(slide, leader_left, y + Inches(0.25), Inches(1.5),
                   "border", Pt(1))

        # ページ番号
        if page:
            tf_pg = make_textbox(slide, MARGIN_LR + Inches(10.8), y,
                                 Inches(0.8), Inches(0.4))
            first_paragraph(tf_pg, str(page), size=SZ_BODY,
                            color_token="subtext", align=PP_ALIGN.RIGHT)

    return slide


def create_section_divider(prs, data, pat):
    """4-3. セクション区切りスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["section_bg"])

    number = data.get("number", "")
    title = data.get("title", "")
    description = data.get("description", "")

    # Clean Pro: 左に primary 縦線
    if pat.get("section_accent"):
        add_line_v(slide, Inches(0.5), Inches(0), Inches(7.5),
                   pat["section_accent"], Pt(4))

    # デコレーション（secondary の小さい四角）
    add_rect(slide, MARGIN_LR + Inches(0.6), Inches(2.0),
             Inches(0.25), Inches(0.25), "secondary")

    # セクション番号
    if number:
        tf_num = make_textbox(slide, MARGIN_LR + Inches(0.6), Inches(2.5),
                              Inches(3.0), Inches(1.0))
        first_paragraph(tf_num, str(number), size=SZ_KPI,
                        color_token="primary", bold=True)

    # タイトル
    title_y = Inches(3.5) if number else Inches(2.8)
    tf_title = make_textbox(slide, MARGIN_LR + Inches(0.6), title_y,
                            Inches(10.0), Inches(1.0))
    first_paragraph(tf_title, title, size=Pt(26),
                    color_token="text", bold=True)

    # 概要テキスト
    if description:
        tf_desc = make_textbox(slide, MARGIN_LR + Inches(0.6), title_y + Inches(1.0),
                               Inches(8.0), Inches(0.8))
        first_paragraph(tf_desc, description, size=SZ_BODY, color_token="subtext")

    return slide


def create_content_text(prs, data, pat):
    """4-4. コンテンツスライド：テキスト主体"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    title_text = data.get("title", "")
    add_slide_header(slide, title_text, pat)

    # テキスト（段落モード）
    text = data.get("text", "")
    # 箇条書き（content または bullets）
    bullets = data.get("content", data.get("bullets", []))

    if text or bullets:
        tf = make_textbox(slide, MARGIN_LR + Inches(0.2), CONTENT_TOP,
                          CONTENT_WIDTH - Inches(0.4), CONTENT_HEIGHT)

        first_set = True
        if text:
            first_paragraph(tf, text, size=SZ_BODY, color_token="text",
                            space_after=Pt(12))
            first_set = False

        for item in bullets:
            if first_set:
                p = tf.paragraphs[0] if not text else tf.add_paragraph()
                if not text:
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(8)
                    add_run(p, item, size=SZ_BODY, color_token="text")
                else:
                    p = add_paragraph(tf, item, size=SZ_BODY,
                                      color_token="text", space_after=Pt(8))
                first_set = False
            else:
                p = add_paragraph(tf, item, size=SZ_BODY,
                                  color_token="text", space_after=Pt(8))

            # 箇条書きマーカー設定
            pPr = p._p.get_or_add_pPr()
            buNone = pPr.find(qn("a:buNone"))
            if buNone is not None:
                pPr.remove(buNone)
            buClr = pPr.makeelement(qn("a:buClr"), {})
            srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "C07A5A"})
            buClr.append(srgb)
            pPr.append(buClr)
            buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u25cf"})
            pPr.append(buChar)
            buSz = pPr.makeelement(qn("a:buSzPct"), {"val": "60000"})
            pPr.append(buSz)
            indent = pPr.makeelement(qn("a:spcBef"), {})
            indent_pts = indent.makeelement(qn("a:spcPts"), {"val": "400"})
            indent.append(indent_pts)
            pPr.append(indent)

    return slide


def create_content_cards(prs, data, pat):
    """4-5. コンテンツスライド：カード型"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    title_text = data.get("title", "")
    add_slide_header(slide, title_text, pat)

    cards = data.get("cards", [])
    n = min(len(cards), 4)
    if n == 0:
        return slide

    gap = Inches(0.2)
    total_gap = gap * (n - 1)
    card_w = (CONTENT_WIDTH - total_gap) / n
    card_h = Inches(3.5)
    card_top = CONTENT_TOP + Inches(0.3)

    card_bg = pat["card_bg"]
    card_border = pat.get("card_border")
    card_radius = pat["card_radius"]

    icon_colors = ["primary", "secondary", "accent", "primary"]

    for i in range(n):
        card = cards[i]
        card_left = MARGIN_LR + (card_w + gap) * i

        # カード背景
        shape = add_rounded_rect(slide, card_left, card_top, card_w, card_h,
                                 card_bg, card_border, card_radius)

        icon_color = icon_colors[i % len(icon_colors)]

        # アイコン（円 or 角丸四角）
        icon_size = Inches(0.35)
        icon_left = card_left + Inches(0.3)
        icon_top = card_top + Inches(0.3)
        if pat["card_radius"] <= 0.03:
            # Clean Pro: 角丸四角アイコン
            icon_shape = add_rounded_rect(slide, icon_left, icon_top,
                                          icon_size, icon_size,
                                          icon_color, radius=0.15)
        else:
            # 円形アイコン
            icon_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, icon_left, icon_top, icon_size, icon_size
            )
            icon_shape.fill.solid()
            icon_shape.fill.fore_color.rgb = C(icon_color)
            icon_shape.line.fill.background()

        # ラベル
        tf_label = make_textbox(slide, card_left + Inches(0.3),
                                card_top + Inches(0.9),
                                card_w - Inches(0.6), Inches(0.5))
        first_paragraph(tf_label, card.get("label", ""),
                        size=SZ_CARD_LABEL, color_token="text", bold=True)

        # 説明文
        desc = card.get("description", "")
        if desc:
            tf_desc = make_textbox(slide, card_left + Inches(0.3),
                                   card_top + Inches(1.4),
                                   card_w - Inches(0.6), Inches(1.8))
            first_paragraph(tf_desc, desc, size=SZ_CARD_DESC,
                            color_token="subtext")

    return slide


def create_content_data(prs, data, pat):
    """4-6. コンテンツスライド：データ・チャート"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    title_text = data.get("title", "")
    add_slide_header(slide, title_text, pat)

    metrics = data.get("metrics", [])
    description = data.get("description", "")
    n = min(len(metrics), 4)

    if n > 0:
        gap = Inches(0.3)
        total_gap = gap * (n - 1)
        box_w = (CONTENT_WIDTH - total_gap) / n
        box_h = Inches(1.8)
        box_top = CONTENT_TOP + Inches(0.2)

        for i in range(n):
            m = metrics[i]
            box_left = MARGIN_LR + (box_w + gap) * i

            # KPI ボックス背景
            add_rounded_rect(slide, box_left, box_top, box_w, box_h,
                             "white", pat.get("card_border"), pat["card_radius"])

            # 数値
            tf_val = make_textbox(slide, box_left, box_top + Inches(0.2),
                                  box_w, Inches(0.9))
            first_paragraph(tf_val, str(m.get("value", "")), size=SZ_KPI,
                            color_token="primary", bold=True, align=PP_ALIGN.CENTER)

            # ラベル
            tf_lbl = make_textbox(slide, box_left, box_top + Inches(1.1),
                                  box_w, Inches(0.4))
            first_paragraph(tf_lbl, m.get("label", ""), size=SZ_CAPTION,
                            color_token="subtext", align=PP_ALIGN.CENTER)

    # 説明テキスト（メトリクス下）
    if description:
        desc_top = CONTENT_TOP + Inches(2.5) if n > 0 else CONTENT_TOP
        tf_desc = make_textbox(slide, MARGIN_LR, desc_top,
                               CONTENT_WIDTH, Inches(2.5))
        first_paragraph(tf_desc, description, size=SZ_BODY, color_token="text")

    return slide


def create_content_table(prs, data, pat):
    """4-7. コンテンツスライド：比較・テーブル"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    title_text = data.get("title", "")
    add_slide_header(slide, title_text, pat)

    headers = data.get("headers", [])
    rows = data.get("rows", [])

    if not headers or not rows:
        return slide

    cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    table_left = MARGIN_LR
    table_top = CONTENT_TOP + Inches(0.2)
    table_width = CONTENT_WIDTH
    row_h = Inches(0.5)
    table_height = row_h * n_rows

    table_shape = slide.shapes.add_table(n_rows, cols,
                                         table_left, table_top,
                                         table_width, table_height)
    table = table_shape.table

    header_bg = pat["table_header_bg"]
    header_text = pat["table_header_text"]

    # ヘッダー行
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, str(h), size=SZ_TABLE_CELL, color_token=header_text, bold=True)
        # ヘッダー背景
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = C(header_bg)

    # データ行
    for i, row in enumerate(rows):
        row_bg = "bg" if i % 2 == 0 else "white"
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = ""
            val = str(row[j]) if j < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            # 強調セル判定
            highlight_cells = data.get("highlight_cells", [])
            is_highlight = [i, j] in highlight_cells or (i, j) in highlight_cells
            if is_highlight:
                add_run(p, val, size=SZ_TABLE_CELL, color_token="primary", bold=True)
            else:
                add_run(p, val, size=SZ_TABLE_CELL, color_token="text")

            # セル背景
            cell_fill = cell.fill
            cell_fill.solid()
            cell_fill.fore_color.rgb = C(row_bg)

    # テーブル枠線
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        tblPr = tbl.makeelement(qn("a:tblPr"), {})
        tbl.insert(0, tblPr)
    # 既存のスタイル参照を削除
    tblStyle = tblPr.find(qn("a:tblStyle"))
    if tblStyle is not None:
        tblPr.remove(tblStyle)
    # バンディングをオフ
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("lastRow", "0")

    return slide


def create_closing_slide(prs, data, pat):
    """4-8. クロージングスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["closing_bg"])

    message = data.get("message", data.get("title", "ありがとうございました"))
    submessage = data.get("submessage", "")
    cta = data.get("cta", "")
    contact = data.get("contact", [])
    heading_color = pat["heading_color"]

    # メインメッセージ
    tf_msg = make_textbox(slide, Inches(1.5), Inches(1.8),
                          Inches(10.333), Inches(1.2))
    first_paragraph(tf_msg, message, size=Pt(26),
                    color_token=heading_color, bold=True, align=PP_ALIGN.CENTER)

    # 区切り線
    line_w = Inches(0.7)
    line_left = (SLIDE_WIDTH - line_w) / 2
    add_line_h(slide, line_left, Inches(3.0), line_w, "primary", Pt(2))

    # サブメッセージ
    if submessage:
        tf_sub = make_textbox(slide, Inches(1.5), Inches(3.3),
                              Inches(10.333), Inches(0.6))
        first_paragraph(tf_sub, submessage, size=SZ_BODY,
                        color_token="subtext", align=PP_ALIGN.CENTER)

    # CTA ボタン
    if cta:
        btn_w = Inches(3.5)
        btn_h = Inches(0.6)
        btn_left = (SLIDE_WIDTH - btn_w) / 2
        btn_top = Inches(4.2)
        btn_shape = add_rounded_rect(slide, btn_left, btn_top, btn_w, btn_h,
                                     pat["cta_bg"], radius=pat["cta_radius"])
        btn_shape.text_frame.word_wrap = True
        p = btn_shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_run(p, cta, size=SZ_BODY, color_token=pat["cta_text"], bold=True)
        # 垂直中央
        btn_shape.text_frame.paragraphs[0].space_before = Pt(4)

    # 連絡先
    if contact:
        contact_top = Inches(5.2) if cta else Inches(4.2)
        tf_contact = make_textbox(slide, Inches(1.5), contact_top,
                                  Inches(10.333), Inches(1.5))
        for i, line in enumerate(contact):
            if i == 0:
                first_paragraph(tf_contact, line, size=SZ_CAPTION,
                                color_token="subtext", align=PP_ALIGN.CENTER,
                                space_after=Pt(4))
            else:
                add_paragraph(tf_contact, line, size=SZ_CAPTION,
                              color_token="subtext", align=PP_ALIGN.CENTER,
                              space_after=Pt(4))

    return slide


# ============================================================
# 後方互換レイアウト
# ============================================================

def create_title_content_slide(prs, data, pat):
    """title_content の後方互換（content_text にマッピング）"""
    return create_content_text(prs, data, pat)


def create_two_column_slide(prs, data, pat):
    """2カラム比較スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    title_text = data.get("title", "")
    add_slide_header(slide, title_text, pat)

    heading_color = pat["heading_color"]

    # 左カラム
    left_data = data.get("left", {})
    left_header = left_data.get("header", "")
    left_items = left_data.get("items", [])

    tf_left = make_textbox(slide, MARGIN_LR, CONTENT_TOP,
                           Inches(5.5), CONTENT_HEIGHT)
    first_paragraph(tf_left, left_header, size=SZ_COLUMN_HEADER,
                    color_token=heading_color, bold=True, space_after=Pt(12))
    for item in left_items:
        add_paragraph(tf_left, f"\u2022  {item}", size=SZ_BODY,
                      color_token="text", space_after=Pt(6))

    # 中央区切り線
    add_line_v(slide, Inches(6.5), CONTENT_TOP, CONTENT_HEIGHT,
               "border", Pt(1))

    # 右カラム
    right_data = data.get("right", {})
    right_header = right_data.get("header", "")
    right_items = right_data.get("items", [])

    tf_right = make_textbox(slide, Inches(7.0), CONTENT_TOP,
                            Inches(5.5), CONTENT_HEIGHT)
    first_paragraph(tf_right, right_header, size=SZ_COLUMN_HEADER,
                    color_token=heading_color, bold=True, space_after=Pt(12))
    for item in right_items:
        add_paragraph(tf_right, f"\u2022  {item}", size=SZ_BODY,
                      color_token="text", space_after=Pt(6))

    return slide


def create_blank_slide(prs, data, pat):
    """白紙スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, pat["slide_bg"])

    title_text = data.get("title", "")
    if title_text:
        add_slide_header(slide, title_text, pat)

    text = data.get("text", "")
    if text:
        top = CONTENT_TOP if title_text else Inches(1.0)
        tf = make_textbox(slide, MARGIN_LR, top, CONTENT_WIDTH, CONTENT_HEIGHT)
        first_paragraph(tf, text, size=SZ_BODY, color_token="text")

    return slide


# ============================================================
# レイアウト登録
# ============================================================

LAYOUT_HANDLERS = {
    # 新レイアウト
    "title_slide": create_title_slide,
    "toc": create_toc_slide,
    "section_divider": create_section_divider,
    "content_text": create_content_text,
    "content_cards": create_content_cards,
    "content_data": create_content_data,
    "content_table": create_content_table,
    "closing": create_closing_slide,
    # 後方互換
    "title_content": create_title_content_slide,
    "two_column": create_two_column_slide,
    "blank": create_blank_slide,
}


# ============================================================
# バリデーション・生成・CLI
# ============================================================

def validate_input(data):
    """入力データのバリデーション"""
    errors = []
    if "title" not in data:
        errors.append("トップレベルに 'title' フィールドが必要です。")
    if "slides" not in data:
        errors.append("トップレベルに 'slides' フィールドが必要です。")
    elif not isinstance(data["slides"], list):
        errors.append("'slides' は配列である必要があります。")
    elif len(data["slides"]) == 0:
        errors.append("'slides' には少なくとも1つのスライドが必要です。")
    else:
        for i, slide in enumerate(data["slides"]):
            if "layout" not in slide:
                errors.append(f"スライド {i + 1}: 'layout' フィールドが必要です。")
            elif slide["layout"] not in LAYOUT_HANDLERS:
                valid = ", ".join(sorted(LAYOUT_HANDLERS.keys()))
                errors.append(
                    f"スライド {i + 1}: 不明なレイアウト '{slide['layout']}'。"
                    f"有効なレイアウト: {valid}"
                )
    return errors


def generate_pptx(input_data, output_path, pattern_override=None):
    """PPTX ファイルを生成する"""
    # パターン解決
    raw_name = (pattern_override
                or input_data.get("pattern")
                or input_data.get("template", DEFAULT_PATTERN))
    pattern_name, pat = resolve_pattern(raw_name)
    print(f"デザインパターン: {pattern_name}")

    # テンプレートファイルの解決
    script_dir = Path(__file__).parent
    skill_dir = script_dir.parent
    template_path = skill_dir / "assets" / "templates" / f"{pattern_name}.pptx"

    if template_path.exists():
        print(f"テンプレート使用: {template_path}")
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    # スライドサイズ 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # タイトルスライドの自動生成
    slides = input_data.get("slides", [])
    first_is_title = slides and slides[0].get("layout") == "title_slide"
    if not first_is_title:
        title_data = {
            "title": input_data.get("title", ""),
            "subtitle": input_data.get("subtitle", ""),
            "author": input_data.get("author", ""),
            "date": input_data.get("date", ""),
        }
        create_title_slide(prs, title_data, pat)

    # 各スライド生成
    for i, slide_data in enumerate(slides):
        layout = slide_data.get("layout", "blank")
        handler = LAYOUT_HANDLERS.get(layout)
        if handler is None:
            print(f"警告: スライド {i + 1}: 不明なレイアウト '{layout}' をスキップ")
            continue

        if layout == "title_slide":
            slide_data.setdefault("subtitle", input_data.get("subtitle", ""))
            slide_data.setdefault("author", input_data.get("author", ""))
            slide_data.setdefault("date", input_data.get("date", ""))

        handler(prs, slide_data, pat)
        print(f"スライド {i + 1}: {layout}")

    # 保存
    output_dir = Path(output_path).parent
    output_dir.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"\n保存: {output_path}（{len(prs.slides)} スライド）")


def main():
    parser = argparse.ArgumentParser(
        description="JSON入力から PowerPoint スライドを生成する"
    )
    parser.add_argument("--input", "-i", required=True, help="入力JSONファイル")
    parser.add_argument("--output", "-o", required=True, help="出力PPTXファイル")
    parser.add_argument(
        "--pattern", "-p", default=None,
        help="デザインパターン（soft_business / clean_pro / warm_formal）"
    )
    # 後方互換: --template も受け付ける
    parser.add_argument("--template", "-t", default=None, help=argparse.SUPPRESS)
    args = parser.parse_args()

    pattern = args.pattern or args.template

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"エラー: 入力ファイルが見つかりません: {args.input}")
        sys.exit(1)

    try:
        with open(input_path, "r", encoding="utf-8") as f:
            input_data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"エラー: JSON解析失敗: {e}")
        sys.exit(1)

    errors = validate_input(input_data)
    if errors:
        print("入力データエラー:")
        for err in errors:
            print(f"  - {err}")
        sys.exit(1)

    generate_pptx(input_data, args.output, pattern_override=pattern)


if __name__ == "__main__":
    main()
